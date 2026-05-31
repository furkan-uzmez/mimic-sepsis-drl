from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path('iql_final_presentation.pptx')
FIG = Path('figures')

prs = Presentation()
prs.slide_width = Inches(13.333333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

BG = RGBColor(7, 17, 31)
PANEL = RGBColor(10, 31, 52)
CYAN = RGBColor(63, 245, 255)
MAGENTA = RGBColor(255, 60, 172)
LIME = RGBColor(217, 255, 88)
INK = RGBColor(232, 251, 255)
MUTED = RGBColor(150, 180, 194)
WHITE = RGBColor(255, 255, 255)

TITLE_FONT = 'Aptos Display'
BODY_FONT = 'Aptos'


def blank():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_grid(slide)
    return slide


def add_grid(slide):
    # Abstract cyber grid / atmosphere.
    for i in range(0, 18):
        x = Inches(i * 0.78)
        line = slide.shapes.add_connector(1, x, 0, x, H)
        line.line.color.rgb = RGBColor(18, 53, 72)
        line.line.transparency = 70
        line.line.width = Pt(0.35)
    for i in range(0, 11):
        y = Inches(i * 0.75)
        line = slide.shapes.add_connector(1, 0, y, W, y)
        line.line.color.rgb = RGBColor(18, 53, 72)
        line.line.transparency = 75
        line.line.width = Pt(0.35)
    orb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.1), Inches(5.05), Inches(2.8), Inches(2.8))
    orb.fill.solid(); orb.fill.fore_color.rgb = MAGENTA; orb.fill.transparency = 70
    orb.line.fill.background()
    orb2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(0.25), Inches(2.3), Inches(2.3))
    orb2.fill.solid(); orb2.fill.fore_color.rgb = CYAN; orb2.fill.transparency = 78
    orb2.line.fill.background()


def text_box(slide, text, x, y, w, h, size=24, color=INK, bold=False, font=BODY_FONT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def eyebrow(slide, text, x=0.75, y=0.55):
    return text_box(slide, text.upper(), x, y, 5.7, 0.35, size=10.5, color=CYAN, bold=True)


def title(slide, text, x=0.75, y=1.2, w=5.7, h=1.5, size=34):
    return text_box(slide, text, x, y, w, h, size=size, color=INK, bold=True, font=TITLE_FONT)


def body(slide, text, x, y, w, h, size=15, color=MUTED):
    return text_box(slide, text, x, y, w, h, size=size, color=color)


def panel(slide, x, y, w, h, accent=CYAN):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = PANEL; shp.fill.transparency = 8
    shp.line.color.rgb = accent; shp.line.transparency = 30; shp.line.width = Pt(1.0)
    return shp


def metric(slide, x, y, w, h, value, label, accent=CYAN):
    panel(slide, x, y, w, h, accent)
    text_box(slide, value, x + 0.18, y + 0.15, w - 0.36, h * 0.45, size=25, color=INK, bold=True, font=TITLE_FONT)
    body(slide, label, x + 0.18, y + h * 0.57, w - 0.36, h * 0.32, size=10.5, color=MUTED)


def bullets(slide, items, x, y, w, h, size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = BODY_FONT
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(8)
        p.bullet = True
    return box


def image_slide(kicker, headline, desc, img, caption):
    s = blank(); eyebrow(s, kicker); title(s, headline, y=1.05, w=4.6, h=1.75, size=29)
    body(s, desc, 0.75, 3.05, 4.35, 1.15, 14)
    panel(s, 5.35, 0.82, 7.25, 5.75, CYAN)
    s.shapes.add_picture(str(FIG / img), Inches(5.55), Inches(1.05), width=Inches(6.85), height=Inches(4.85))
    body(s, caption, 5.55, 6.05, 6.85, 0.35, 10.5)


# 1
s = blank(); eyebrow(s, 'IEEE final report sunumu'); title(s, 'Destek sınırlı sepsis dozaj politikaları', y=1.15, w=5.9, h=1.9, size=36)
body(s, 'IQL tabanlı, sızıntısız ve denetlenebilir bir çevrimdışı pekiştirmeli öğrenme iş akışı.', 0.75, 3.35, 5.45, 0.85, 15.5, RGBColor(198, 233, 242))
metric(s, 7.0, 1.15, 2.2, 1.15, '72s', 'akut gözlem penceresi')
metric(s, 9.45, 1.15, 2.2, 1.15, '62', 'boyutlu durum vektörü', MAGENTA)
metric(s, 7.0, 2.65, 2.2, 1.15, '25', 'sıvı-vazopresör aksiyonu', LIME)
metric(s, 9.45, 2.65, 2.2, 1.15, '30', 'final eğitim koşusu')
body(s, 'MIMIC-IV v3.1 / Offline RL', 0.75, 6.75, 3.4, 0.25, 10, CYAN)

# 2
s = blank(); eyebrow(s, 'Problem'); title(s, 'Gözlemsel veride iyi görünen politika güvenli olmayabilir.', y=1.05, w=5.2, h=2.0, size=30)
body(s, 'Sepsiste agresif müdahale genellikle daha ağır hastalarda gözlenir; değer tahmini destek eksikliğiyle kolayca yanılır.', 0.75, 3.3, 4.9, 1.0, 14.5)
items = [('Ekstrapolasyon', 'Zayıf temsil edilen doz kombinasyonlarına yapay değer atanabilir.'), ('OPE varyansı', 'WIS birkaç yüksek ağırlıklı yörüngeye yaslanabilir.'), ('Sızıntı riski', 'Ön işleme veya split hatası final metriği şişirir.'), ('Klinik sınır', 'Çalışma yatak başı öneri değil, retrospektif protokoldür.')]
for i, (h, d) in enumerate(items):
    x = 6.35 + (i % 2) * 3.0; y = 1.05 + (i // 2) * 2.0
    panel(s, x, y, 2.65, 1.55, CYAN if i % 2 == 0 else MAGENTA)
    text_box(s, h, x + .18, y + .17, 2.3, .35, 16, INK, True, TITLE_FONT)
    body(s, d, x + .18, y + .65, 2.25, .65, 11.5)

# 3
s = blank(); eyebrow(s, 'Uçtan uca mimari'); title(s, 'Tek model değil, denetlenebilir deney hattı.', y=.95, w=7.5, h=.9, size=30)
steps = ['Sepsis-3 kohort', '72 saat / 4 saatlik karar', 'Hasta düzeyinde split', '62 özellikli durum', '25 aksiyon + ödül', 'Replay + baseline', 'IQL eğitimi', '18 aday tarama', 'OPE ve güvenlik', 'Çoklu tohum final']
for i, st in enumerate(steps):
    x = .75 + (i % 5) * 2.45; y = 2.15 + (i // 5) * 1.75
    panel(s, x, y, 2.15, 1.25, CYAN if i < 5 else MAGENTA)
    text_box(s, f'{i+1:02d}', x + .16, y + .12, .7, .32, 18, CYAN, True, TITLE_FONT)
    body(s, st, x + .16, y + .55, 1.8, .45, 11.3, INK)

# 4
s = blank(); eyebrow(s, 'MDP formülasyonu'); title(s, 'Karar birimi: 4 saatlik klinik pencere.', y=1.0, w=4.8, h=1.6, size=31)
panel(s, 6.0, .95, 6.45, 5.1, CYAN)
bullets(s, ['Durum: vital, laboratuvar, tedavi geçmişi, demografi ve missingness sinyalleri.', 'Aksiyon: vazopresör ve IV sıvı binlerinin 5 x 5 ayrık kombinasyonu.', 'Ödül: terminal sağkalım/ölüm faydası ve opsiyonel SOFA şekillendirmesi.', 'İskonto: gamma = 0.99 ile uzun vadeli mortalite etkisi korunur.', 'Sızıntı önlemi: binler ve ön işleme yalnız eğitim bölmesinde fit edilir.'], 6.35, 1.35, 5.7, 4.2, 14)

image_slide('Aşama 1', 'Değer tek başına seçici değil; destekle birlikte okunur.', 'SOFA-konservatif adaylar öne çıktı; final havuzu tek metriğe aşırı uyumu engellemek için çeşitlendirildi.', 'fqe_vs_support.png', 'Yüksek FQE ancak yeterli davranış desteğiyle anlamlı kabul edilir.')

# 6
s = blank(); eyebrow(s, 'IQL tarama tasarımı'); title(s, '18 adaydan 6 finalist, sonra 3 tohumla tekrar ölçüm.', y=1.0, w=6.0, h=1.4, size=30)
cols = [('Ödül', 'Sparse / SOFA', 'Terminal mortalite sinyali SOFA değişimiyle veya yalnız başına değerlendirilir.'), ('LR rejimi', 'Konservatif / referans / aktör-konservatif', 'Eleştirmen ve aktör agresifliği destek sınırları açısından ayrıştırılır.'), ('IQL ayarı', 'Güvenli / referans / iyimser', 'Ekspektil ve sıcaklık politika çıkarımının risk profilini kontrol eder.')]
for i, (k, h, d) in enumerate(cols):
    x = .9 + i * 4.1
    panel(s, x, 3.05, 3.55, 2.3, [CYAN, MAGENTA, LIME][i])
    text_box(s, k.upper(), x + .22, 3.28, 2.8, .25, 10, [CYAN, MAGENTA, LIME][i], True)
    text_box(s, h, x + .22, 3.72, 3.0, .65, 18, INK, True, TITLE_FONT)
    body(s, d, x + .22, 4.55, 3.05, .55, 11.8)

image_slide('Aşama 2', 'Final seçim Pareto ödünleşiminden gelir.', 'Seçili yapı: iql_sofa_shaped_conservative_safe. Karar FQE/WIS değerlerini ESS, destek ve klinisyen uyumuyla birlikte tartar.', 'pareto_frontier.png', 'Değer, destek ve güvenlik eksenleri birlikte final kararı belirler.')

# 8
s = blank(); eyebrow(s, 'Final test sonuçları'); title(s, 'Metodolojik kanıt güçlü; klinik üstünlük iddiası yok.', y=1.0, w=5.1, h=1.6, size=30)
body(s, 'ESS 50 eşiğinin altında olduğu için sonuçlar prospektif performans kanıtı olarak değil, destek sınırları tanımlı OPE kanıtı olarak yorumlanır.', .75, 3.0, 4.8, 1.0, 14)
metrics = [('2.848','FQE'),('8.203','WIS'),('29.41','ESS'),('0.991','davranışsal destek'),('0.414','klinisyen tam-kutu uyumu'),('4.963-10.817','WIS %95 CI')]
for i,(v,l) in enumerate(metrics):
    x = 6.15 + (i % 2) * 2.9; y = .9 + (i // 2) * 1.55
    metric(s, x, y, 2.55, 1.12, v, l, [CYAN, MAGENTA, LIME][i % 3])

image_slide('Baseline karşılaştırması', 'Yüksek WIS, düşük ESS ile güvenilir değildir.', 'Davranış klonlama WIS açısından yüksek görünür, fakat ESS=1.6 düzeyi nedeniyle istatistiksel olarak kırılgandır.', 'baseline_comparison.png', 'Klinisyen replay, no-treatment, davranış klonlama ve seçili IQL birlikte okunur.')
image_slide('Aksiyon güvenliği', 'Politika hekimden sapıyor, fakat veri desteği içinde kalıyor.', 'Desteklenmeyen aksiyon sapması 0.009 düzeyinde; farklı kararlar çoğunlukla gözlenmiş makul komşu kutularda kalır.', 'action_heatmap.png', '5 x 5 tedavi gridi sıvı ve vazopresör şiddetini birlikte kodlar.')

# 11
s = blank(); eyebrow(s, 'Sınırlar ve etik'); title(s, 'Retrospektif model, klinik karar destek sistemi değildir.', y=1.0, w=5.4, h=1.5, size=30)
panel(s, 6.1, .95, 6.2, 4.9, MAGENTA)
bullets(s, ['Durum vektörü hekim niyetini, kontrendikasyonları ve tüm yatak başı gözlemleri kapsamaz.', 'MIMIC-IV kullanımı credentialing, CITI eğitimi ve veri kullanım sözleşmesi yükümlülüklerine tabidir.', 'FQE/WIS prospektif klinik kanıt değil; model ve davranış politikası varsayımlarına bağlıdır.', 'Düşük ESS, bootstrap aralığı ve gözlemsel confounding temkinli raporlama gerektirir.'], 6.45, 1.45, 5.45, 3.8, 14)

# 12
s = blank(); eyebrow(s, 'Gelecek yönler'); title(s, 'Sonraki adım daha büyük grid değil, belirsizlik duyarlı öğrenme.', y=.95, w=5.5, h=1.7, size=29)
items = [('Dinamik ekspektil', 'Yoğun destek bölgelerinde agresif, seyrek bölgelerde davranışa yakın politika çıkarımı.'), ('Multi-step izler', 'Seyrek terminal mortalite sinyalini daha hızlı ve destek-duyarlı yayma.'), ('Manifold düzenleme', 'Gerçekçi sınır durumlarıyla değer fonksiyonuna güvenli marjlar öğretme.'), ('Risk modelleri', 'Quantile tabanlı belirsizlikle yüksek riskli outlier etkisini sınırlama.')]
for i, (h, d) in enumerate(items):
    x = 6.25 + (i % 2) * 3.0; y = 1.0 + (i // 2) * 2.05
    panel(s, x, y, 2.65, 1.58, CYAN if i % 2 == 0 else LIME)
    text_box(s, h, x + .18, y + .18, 2.25, .35, 15.5, INK, True, TITLE_FONT)
    body(s, d, x + .18, y + .68, 2.25, .62, 11.5)

# 13
s = blank(); eyebrow(s, 'Sonuç'); title(s, 'IQL hattı güvenli sınırları görünür kılan bir değerlendirme protokolü sunar.', y=.95, w=6.0, h=1.8, size=29)
body(s, 'Kohorttan rapora kadar aynı protokolde izlenen süreç, sepsis dozaj politikalarını denetlenebilir ve tekrar üretilebilir biçimde karşılaştırmayı sağlar.', .75, 3.1, 5.3, .95, 14)
panel(s, 6.7, 1.1, 5.35, 4.3, CYAN)
bullets(s, ['Aşama 0 denetimi sızıntı ve sözleşme kontrollerini geçti.', 'Aşama 1 destek-duyarlı 6 finalist üretti.', 'Aşama 2 final konfigürasyonu çoklu tohumla doğruladı.', 'Nihai yorum: retrospektif OPE kanıtı, klinik üstünlük kanıtı değil.'], 7.05, 1.65, 4.65, 3.1, 14)

# Slide numbers
for idx, slide in enumerate(prs.slides, start=1):
    text_box(slide, f'{idx:02d} / {len(prs.slides):02d}', 11.55, 6.92, 1.1, .25, 9.5, MUTED, True, align=PP_ALIGN.RIGHT)

prs.save(OUT)
print(OUT.resolve())
