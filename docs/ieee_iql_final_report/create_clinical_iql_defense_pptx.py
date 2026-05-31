from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path("clinical_iql_defense_presentation.pptx")
FIG = Path("figures")

prs = Presentation()
prs.slide_width = Inches(13.333333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height

PAPER = RGBColor(244, 239, 227)
INK = RGBColor(17, 21, 28)
MUTED = RGBColor(88, 96, 109)
RED = RGBColor(229, 54, 46)
TEAL = RGBColor(11, 102, 112)
GOLD = RGBColor(208, 155, 44)
CREAM = RGBColor(255, 250, 242)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(31, 38, 50)

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"


def add_text(slide, text, x, y, w, h, size=20, color=INK, bold=False,
             font=BODY_FONT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    # Left tri-color clinical spine.
    for y, col in [(0, RED), (2.5, TEAL), (5.0, GOLD)]:
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(y), Inches(0.16), Inches(2.5))
        s.fill.solid(); s.fill.fore_color.rgb = col; s.line.fill.background()
    # Quiet dossier grid.
    for i in range(0, 18):
        x = Inches(0.55 + i * 0.75)
        line = slide.shapes.add_connector(1, x, 0, x, H)
        line.line.color.rgb = RGBColor(214, 205, 188)
        line.line.transparency = 55
        line.line.width = Pt(0.25)
    for i in range(0, 11):
        y = Inches(0.3 + i * 0.7)
        line = slide.shapes.add_connector(1, 0, y, W, y)
        line.line.color.rgb = RGBColor(214, 205, 188)
        line.line.transparency = 65
        line.line.width = Pt(0.25)
    # Faint corner circles for atmosphere.
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.9), Inches(-0.55), Inches(2.4), Inches(2.4))
    c1.fill.solid(); c1.fill.fore_color.rgb = TEAL; c1.fill.transparency = 82; c1.line.fill.background()
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.55), Inches(5.65), Inches(3.3), Inches(3.3))
    c2.fill.solid(); c2.fill.fore_color.rgb = RED; c2.fill.transparency = 88; c2.line.fill.background()


def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    return s


def panel(slide, x, y, w, h, fill=CREAM, border=RGBColor(204, 193, 174), radius=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.fill.transparency = 2
    shp.line.color.rgb = border; shp.line.transparency = 18; shp.line.width = Pt(1.0)
    return shp


def eyebrow(slide, text, x=0.62, y=0.48, color=TEAL):
    return add_text(slide, text.upper(), x, y, 6.1, 0.28, size=9.6, color=color, bold=True)


def title(slide, text, x=0.62, y=0.95, w=5.6, h=1.25, size=29):
    return add_text(slide, text, x, y, w, h, size=size, color=INK, bold=True, font=TITLE_FONT)


def subtitle(slide, text, x=0.62, y=2.55, w=5.25, h=0.88, size=13.2):
    return add_text(slide, text, x, y, w, h, size=size, color=MUTED)


def metric(slide, x, y, w, h, value, label, accent=RED):
    panel(slide, x, y, w, h, WHITE, accent)
    add_text(slide, value, x + 0.16, y + 0.12, w - 0.32, h * 0.45, size=23, color=INK, bold=True, font=TITLE_FONT)
    add_text(slide, label, x + 0.16, y + h * 0.58, w - 0.32, h * 0.35, size=9.8, color=MUTED, bold=True)


def bullets(slide, items, x, y, w, h, size=13.2, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = BODY_FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
        p.bullet = True
    return box


def label(slide, text, x, y, w, color=RED):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.36))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
    add_text(slide, text.upper(), x + 0.12, y + 0.08, w - 0.24, 0.18, size=7.8, color=WHITE, bold=True)


def add_picture_fit(slide, path, x, y, w, h):
    panel(slide, x, y, w, h, WHITE, TEAL)
    if Path(path).exists():
        slide.shapes.add_picture(str(path), Inches(x + 0.14), Inches(y + 0.16), width=Inches(w - 0.28), height=Inches(h - 0.48))
    else:
        add_text(slide, f"Missing figure: {path}", x + 0.3, y + 0.4, w - 0.6, h - 0.8, size=14, color=RED, bold=True)


def image_slide(kicker, headline, text, image, caption, accent=TEAL):
    s = blank(); eyebrow(s, kicker, color=accent); title(s, headline, y=0.92, w=4.75, h=1.8, size=26)
    subtitle(s, text, y=2.92, w=4.55, h=1.25, size=12.4)
    add_picture_fit(s, FIG / image, 5.55, 0.72, 7.15, 5.55)
    add_text(s, caption, 5.72, 6.38, 6.8, 0.35, size=9.4, color=MUTED, bold=True)
    return s


# 1 Title
s = blank(); eyebrow(s, "Akademik savunma destesi", color=RED)
title(s, "Sepsis tedavisinde güvenli çevrimdışı IQL", y=1.05, w=6.1, h=1.8, size=33)
subtitle(s, "MIMIC-IV v3.1 kohortu uzerinde veri sızıntısı kapalı, destek duyarlı ve etik sınırları açık bir retrospektif politika değerlendirme protokolü.", y=3.05, w=5.75, h=1.0, size=14)
for i, (v, l, c) in enumerate([("72h", "akut faz", RED), ("4h", "karar adımı", TEAL), ("62", "durum özelliği", GOLD), ("25", "ortak doz aksiyonu", RED)]):
    metric(s, 7.05 + (i % 2) * 2.45, 1.2 + (i // 2) * 1.52, 2.15, 1.05, v, l, c)
label(s, "IQL + OPE + support diagnostics", 7.05, 4.55, 3.45, TEAL)
add_text(s, "Final mesaj: model bir bedside öneri motoru değil; klinik RL iddialarını güvenlik sınırları içinde denetleyen bir kanıt hattıdır.", 7.05, 5.12, 4.9, 0.8, 13.2, DARK, True)

# 2 Clinical problem
s = blank(); eyebrow(s, "Klinik problem", color=RED); title(s, "Gözlemsel veri yüksek riskli bir ayna etkisi üretir.", w=5.7, h=1.55, size=29)
subtitle(s, "Sepsiste sıvı ve vazopresör titrasyonu dar terapötik aralıkta yapılır; aşırı veya yetersiz resüsitasyon doğrudan organ yetmezliği riskini değiştirir.", y=2.75, w=5.4)
for i, (h, d, c) in enumerate([
    ("Endikasyon karıştırıcılığı", "Ağır hastalar daha agresif tedavi alır; retrospektif korelasyon nedensellik gibi görünebilir.", RED),
    ("Ekstrapolasyon hatası", "Veride bulunmayan doz kombinasyonları bootstrapped Q-değerleriyle şişebilir.", TEAL),
    ("OPE varyansı", "WIS hedef politika saptıkça birkaç yüksek ağırlıklı yörüngeye dayanabilir.", GOLD),
    ("Klinik sınır", "Çıktı optimizasyon önerisi değil, güvenli retrospektif denetimdir.", RED),
]):
    x = 6.45 + (i % 2) * 3.0; y = 1.0 + (i // 2) * 2.05
    panel(s, x, y, 2.65, 1.55, WHITE, c)
    add_text(s, h, x + .18, y + .16, 2.25, .38, 14.2, INK, True, TITLE_FONT)
    add_text(s, d, x + .18, y + .64, 2.25, .68, 10.5, MUTED)

# 3 Cohort and MDP
s = blank(); eyebrow(s, "Kohort ve MDP", color=TEAL); title(s, "Karar birimi: 72 saatlik akut fazda 4 saatlik klinik pencere.", w=6.3, h=1.5, size=28)
panel(s, 6.55, .95, 5.7, 4.9, CREAM, TEAL)
bullets(s, [
    "MIMIC-IV v3.1 içinden Sepsis-3 yetişkin yoğun bakım kohortu çıkarıldı.",
    "Durum vektörü vital bulgular, laboratuvarlar, demografi, tedavi geçmişi ve missingness maskelerini kapsar.",
    "Aksiyon uzayı IV sıvı ve vazopresör doz quartile'larının 5 x 5 ortak grididir.",
    "Ödül terminal sağkalım sinyalini ve SOFA organ yetmezliği geçişlerini birlikte okuyabilir.",
    "Preprocessing ve bin eşikleri yalnızca eğitim bölmesinde fit edilerek geleceğe bakış engellenir.",
], 6.9, 1.36, 4.95, 3.95, 12.2)
for i, (v, l, c) in enumerate([("18", "tipik karar adımı", RED), ("5x5", "doz gridi", TEAL), ("0.99", "iskonto", GOLD)]):
    metric(s, .82 + i * 1.8, 3.88, 1.48, 1.0, v, l, c)

# 4 Workflow
s = blank(); eyebrow(s, "Sızıntısız deney hattı", color=TEAL); title(s, "Tek model değil, izlenebilir 10 adımlı protokol.", w=7.2, size=28)
steps = ["Kohort", "72h pencereler", "Hasta split", "Train-only scaling", "Aksiyon binleri", "Replay set", "IQL eğitim", "18 aday", "OPE teşhis", "3 seed final"]
for i, st in enumerate(steps):
    x = .75 + (i % 5) * 2.45; y = 2.0 + (i // 5) * 1.75
    panel(s, x, y, 2.05, 1.12, WHITE, [RED, TEAL, GOLD, TEAL, RED][i % 5])
    add_text(s, f"{i+1:02d}", x + .15, y + .12, .65, .28, 16, [RED, TEAL, GOLD, TEAL, RED][i % 5], True, TITLE_FONT)
    add_text(s, st, x + .15, y + .53, 1.62, .35, 10.7, INK, True)
add_text(s, "Kontrol ilkesi", .82, 5.78, 1.35, .25, 10, RED, True)
add_text(s, "Aynı hastanın farklı zaman pencereleri train/test arasında paylaşılmaz; validasyon ve test hiçbir ön işleme parametresini tekrar fit etmez.", 2.05, 5.72, 8.7, .45, 12, MUTED)

# 5 IQL method
s = blank(); eyebrow(s, "Algoritmik seçim", color=RED); title(s, "IQL, görülmeyen aksiyonları maksimize etmeden değer öğrenir.", w=5.85, h=1.55, size=28)
subtitle(s, "CQL aşırı ceza ile nadir fakat hayat kurtarıcı rescue therapy örneklerini bastırabilir; IQL destek dışına çıkmayı reddeden daha uygun bir klinik takas sunar.", y=2.85, w=5.4)
panel(s, 6.45, 1.05, 5.5, 1.22, DARK, RED)
add_text(s, "L2^tau(u) = |tau - 1(u < 0)| u^2", 6.72, 1.42, 4.95, .35, 22, WHITE, True, TITLE_FONT, align=PP_ALIGN.CENTER)
for i, (h, d, c) in enumerate([
    ("Expectile", "Başarılı veri içi kararların üst sınırını hedefler.", RED),
    ("Double critics", "Bootstrap önyargısını çift eleştirmenle sınırlar.", TEAL),
    ("Temperature", "Aktörün klinisyenden radikal sapmasını kontrol eder.", GOLD),
]):
    panel(s, 6.55 + i * 1.82, 3.22, 1.55, 1.65, WHITE, c)
    add_text(s, h, 6.70 + i * 1.82, 3.43, 1.2, .28, 12.4, INK, True, TITLE_FONT, align=PP_ALIGN.CENTER)
    add_text(s, d, 6.70 + i * 1.82, 3.88, 1.2, .66, 8.8, MUTED, align=PP_ALIGN.CENTER)

# 6 Hyperparameter search
s = blank(); eyebrow(s, "Tarama tasarımı", color=GOLD); title(s, "18 adaydan 6 finalist: seçim tek metrikle yapılmadı.", w=6.4, size=28)
for i, (h, sub, d, c) in enumerate([
    ("Aşama 1", "18 aday", "Ödül tipi, LR rejimi, ekspektil ve aktör sıcaklığı sistematik tarandı.", RED),
    ("Aşama 2", "6 finalist", "Yüksek değer ile davranış desteği birlikte korunacak şekilde havuz daraltıldı.", TEAL),
    ("Final", "3 seed", "Yakınsama gürültüsünü ayırmak için en iyi adaylar bağımsız tohumlarla tekrarlandı.", GOLD),
]):
    x = .95 + i * 4.05
    panel(s, x, 2.35, 3.45, 2.38, WHITE, c)
    add_text(s, h.upper(), x + .22, 2.62, 2.9, .25, 9.2, c, True)
    add_text(s, sub, x + .22, 3.05, 2.9, .45, 24, INK, True, TITLE_FONT)
    add_text(s, d, x + .22, 3.78, 2.92, .58, 10.7, MUTED)
add_text(s, "Nihai aday: iql_sofa_shaped_conservative_safe", .95, 5.55, 5.9, .32, 15, INK, True, TITLE_FONT)
add_text(s, "Seçim kriteri FQE/WIS değerini ESS, support mass, low-support action rate ve klinisyen uyumuyla birlikte tartar.", .95, 6.02, 8.4, .32, 11.5, MUTED)

image_slide("Davranış desteği", "Yüksek FQE ancak destekle anlamlıdır.", "FQE vs support grafiği, düşük destekle şişen değer tahminlerini erken aşamada dışarıda bırakmak için kullanıldı.", "fqe_vs_support.png", "Aşama 2 adayları: değer ve support mass birlikte okunur.", TEAL)
image_slide("Pareto final", "Final seçim değer-risk ödünleşiminden gelir.", "Pareto görünümü final adayın FQE-style skorunu yükseltirken düşük destekli aksiyon oranını yaklaşık %0.9 bandında tuttuğunu gösterir.", "pareto_frontier.png", "Run #46/e200, güvenlik ve değer eksenlerinde Pareto sınırında konumlanır.", RED)

# 9 Results
s = blank(); eyebrow(s, "Nihai test karnesi", color=RED); title(s, "Sonuçlar umut verici; klinik üstünlük kanıtı değil.", w=5.4, h=1.5, size=28)
subtitle(s, "ESS eşiğin altında kaldığı için raporlama prospektif performans iddiası yerine sınırlı ama denetlenebilir OPE kanıtı olarak kurulmalıdır.", y=2.8, w=5.0)
for i, (v, l, c) in enumerate([
    ("2.848", "FQE", RED), ("8.203", "WIS", TEAL), ("29.41", "ESS", GOLD),
    ("0.991", "davranış desteği", TEAL), ("41.4%", "klinisyen uyumu", RED), ("4.963-10.817", "WIS %95 CI", GOLD),
]):
    metric(s, 6.2 + (i % 2) * 2.75, 0.95 + (i // 2) * 1.48, 2.42, 1.02, v, l, c)
label(s, "Dürüst sınırlama", .72, 4.68, 1.65, RED)
add_text(s, "ESS=29.41 değeri 50 güvenlik eşiğinin altında; bu nedenle belirsizlik ve bootstrap CI sunumun merkezi argümanlarından biridir.", .72, 5.22, 4.95, .7, 12.5, INK, True)

image_slide("Baseline kontrol", "Yüksek WIS, düşük ESS ile sahte güven üretir.", "Behavior Cloning güçlü görünür fakat ESS=1.6 düzeyi bu tahminin birkaç uç yörüngeye dayandığını gösterir. Seçili IQL daha dengeli bir varyans profili sunar.", "baseline_comparison.png", "Referans politikalar WIS ve ESS eksenlerinde birlikte karşılaştırılır.", GOLD)
image_slide("Bootstrap", "Geniş aralık metodolojik temkini zorunlu kılar.", "Bootstrap güven aralığı düşük ESS kaynaklı belirsizliği gizlemek yerine görünür kılar; akademik savunmada ana güvenlik bariyeridir.", "bootstrap_ci.png", "WIS belirsizliği sonuçların prospektif iddia olarak okunmasını engeller.", RED)
image_slide("Aksiyon güvenliği", "Politika farklılaşır, ama destek dışına kaçmaz.", "Tam-kutu klinisyen uyumu %41.4; sapmalarda desteklenmeyen aksiyon oranı yalnızca %0.9 seviyesindedir.", "action_heatmap.png", "5 x 5 sıvı-vazopresör gridi politika davranışını okunabilir kılar.", TEAL)
image_slide("Alt grup emniyeti", "Klinik güvenlik alt gruplarda da denetlenir.", "Ciddiyet alt gruplarında destek dışı aksiyonların sınırlı kalması, modelin pasifist veya agresif uçlara kaçmadığını test eder.", "subgroup_safety.png", "SOFA ve ciddiyet katmanlarında güvenlik teşhisleri.", GOLD)

# 14 Limitations
s = blank(); eyebrow(s, "Kısıtlar ve etik", color=RED); title(s, "Bu çalışma klinik karar destek sistemi olarak dağıtılamaz.", w=5.7, size=28)
panel(s, 6.35, .9, 5.85, 4.98, CREAM, RED)
bullets(s, [
    "EHR kayıtları hekimin yatak başı sezgisini, kontrendikasyonları ve tüm klinik nüansları kapsamaz.",
    "Endikasyon bazlı karıştırıcılık gözlemsel kohortta tamamen ortadan kaldırılamaz.",
    "FQE/WIS model varsayımlarına bağlı tahminleyicilerdir; randomize prospektif çalışmanın yerine geçmez.",
    "MIMIC-IV kullanımı PhysioNet credentialing, CITI eğitimi ve HIPAA veri yönetişimi sınırlarıyla yürütülür.",
], 6.72, 1.38, 5.08, 3.85, 12.5)
add_text(s, "Savunma cümlesi", .78, 4.35, 2.6, .3, 14, RED, True, TITLE_FONT)
add_text(s, "Modelin değeri tedavi buyruğu vermesinde değil; hangi RL iddiasının veri desteği içinde kaldığını şeffaf biçimde göstermesindedir.", .78, 4.86, 4.75, .84, 13, INK, True)

# 15 Defensive Q&A
s = blank(); eyebrow(s, "Defansif Q&A", color=TEAL); title(s, "Jürinin sert sorularına hazır üç cevap.", w=6.0, size=29)
qa = [
    ("ESS neden düşük?", "Gizlenmedi; BC'nin ESS=1.6 kırılganlığına karşı IQL 29.41 ile daha stabil, fakat yine de temkinli raporlandı."),
    ("Aksiyon gridi kaba değil mi?", "Evet, fakat continuous OOD ekstrapolasyon riskini azaltan muhafazakar klinik zırh olarak seçildi."),
    ("POMDP ve confounding?", "Missingness sinyalleri, alt grup analizleri ve destek teşhisleri riski azaltır; nedensel kanıt iddiası kurulmaz."),
]
for i, (q, a) in enumerate(qa):
    panel(s, .95 + i * 4.02, 2.25, 3.42, 2.25, WHITE, [RED, TEAL, GOLD][i])
    add_text(s, q, 1.2 + i * 4.02, 2.58, 2.92, .38, 16, INK, True, TITLE_FONT)
    add_text(s, a, 1.2 + i * 4.02, 3.15, 2.78, .88, 10.8, MUTED)
add_text(s, "Sunum tonu: iddialı değil, denetlenebilir ve sınırlarını bilen.", .95, 5.55, 7.5, .32, 15, TEAL, True, TITLE_FONT)

# 16 Future and close
s = blank(); eyebrow(s, "Sonuç ve gelecek çalışma", color=GOLD); title(s, "Güvenli klinik RL, daha büyük aksiyon gridiyle değil belirsizlik katmanıyla ilerler.", w=6.45, h=1.75, size=27)
subtitle(s, "Bu deste, IQL hattını agresif performans anlatısından çıkarıp istatistiksel sınırları görünür kılan bir akademik savunma anlatısına dönüştürür.", y=3.05, w=5.65)
for i, (h, d, c) in enumerate([
    ("Dinamik expectile", "Seyrek bölgede davranışa yakın, yoğun destekte daha esnek politika.", RED),
    ("Uncertainty critics", "Epistemik belirsizlik yüksek durumlarda öneriyi bayraklama.", TEAL),
    ("Guideline filter", "Klinik kılavuz ihlallerini runtime katmanında engelleme.", GOLD),
    ("Prospektif yol", "Simülasyon ve klinisyen-in-the-loop validasyon sonrası sınırlı pilot.", RED),
]):
    x = 6.72 + (i % 2) * 2.75; y = 1.05 + (i // 2) * 2.0
    panel(s, x, y, 2.38, 1.48, WHITE, c)
    add_text(s, h, x + .18, y + .17, 1.98, .32, 13.2, INK, True, TITLE_FONT)
    add_text(s, d, x + .18, y + .64, 1.95, .58, 9.6, MUTED)
add_text(s, "Teşekkürler", .82, 6.17, 3.1, .42, 24, RED, True, TITLE_FONT)

# Slide numbers and footer
for idx, slide in enumerate(prs.slides, start=1):
    add_text(slide, f"{idx:02d} / {len(prs.slides):02d}", 11.52, 6.92, 1.05, .22, size=8.5, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, "MIMIC-Sepsis Offline IQL", .62, 6.92, 2.8, .22, size=8.5, color=MUTED, bold=True)

prs.save(OUT)
print(OUT.resolve())
